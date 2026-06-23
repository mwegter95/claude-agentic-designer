"""Reusable python-pptx helpers: theme/token reading and placeholder filling."""
from __future__ import annotations

from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Emu


def emu_to_pt(value: Optional[int]) -> Optional[float]:
    if value is None:
        return None
    return round(Emu(value).pt, 2)


def layout_catalog(prs: Presentation) -> List[Dict]:
    """Describe every slide layout and its placeholders in semantic terms."""
    catalog: List[Dict] = []
    for li, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append(
                {
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                    "left_pt": emu_to_pt(ph.left),
                    "top_pt": emu_to_pt(ph.top),
                    "width_pt": emu_to_pt(ph.width),
                    "height_pt": emu_to_pt(ph.height),
                }
            )
        catalog.append(
            {
                "index": li,
                "name": layout.name,
                "placeholders": placeholders,
            }
        )
    return catalog


def slide_size(prs: Presentation) -> Tuple[float, float]:
    return emu_to_pt(prs.slide_width), emu_to_pt(prs.slide_height)


def theme_colors(prs: Presentation) -> Dict[str, str]:
    """Extract the theme color map (dk1, lt1, accent1..6, etc.) as hex."""
    colors: Dict[str, str] = {}
    try:
        theme = prs.slide_masters[0].element.getroottree()
    except Exception:
        theme = None
    # python-pptx does not expose theme colors directly; parse the XML part.
    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        from lxml import etree  # python-pptx already depends on lxml

        root = etree.fromstring(theme_part.blob)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        scheme = root.find(".//a:clrScheme", ns)
        if scheme is not None:
            for child in scheme:
                tag = etree.QName(child).localname
                srgb = child.find("a:srgbClr", ns)
                sys = child.find("a:sysClr", ns)
                if srgb is not None:
                    colors[tag] = "#" + srgb.get("val", "").upper()
                elif sys is not None:
                    colors[tag] = "#" + (sys.get("lastClr") or "").upper()
    except Exception:
        pass
    return colors


def theme_fonts(prs: Presentation) -> Dict[str, str]:
    fonts: Dict[str, str] = {}
    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        from lxml import etree

        root = etree.fromstring(theme_part.blob)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        for kind in ("majorFont", "minorFont"):
            el = root.find(f".//a:{kind}/a:latin", ns)
            if el is not None:
                fonts[kind] = el.get("typeface", "")
    except Exception:
        pass
    return fonts
