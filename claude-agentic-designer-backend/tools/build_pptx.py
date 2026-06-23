#!/usr/bin/env python3
"""build_pptx.py — Deterministically render a slide-plan JSON into a .pptx.

The LLM agents NEVER free-draw shapes. They emit a structured slide plan; this
script fills the master's real layouts/placeholders so output stays on-brand and
reproducible. If a master deck is provided, its layouts are used directly.

slide_plan.json shape (see skill/schemas/slide_plan.schema.json):
{
  "doc_type": "slides" | "one_pager" | "white_paper",
  "title": "...",
  "master": "/abs/path/to/master.pptx",   # optional
  "slides": [
    {
      "layout": 1,                          # layout index OR name string
      "placeholders": { "0": "Title text", "1": "Body text" },
      "bullets": { "1": ["point a", "point b"] },   # optional bullet lists
      "notes": "speaker notes"
    }
  ]
}

Usage:
  python tools/build_pptx.py --plan run/slide_plan.json --out run/output.pptx \
      [--run RUN123]
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Optional

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx import Presentation  # noqa: E402
from pptx.util import Pt  # noqa: E402

from tools.common.events import emit  # noqa: E402


def _resolve_layout(prs: Presentation, ref):
    layouts = list(prs.slide_layouts)
    if isinstance(ref, int):
        if 0 <= ref < len(layouts):
            return layouts[ref]
        return layouts[min(ref, len(layouts) - 1)]
    # name match (case-insensitive, fall back to first)
    for layout in layouts:
        if layout.name.strip().lower() == str(ref).strip().lower():
            return layout
    return layouts[0]


def _set_placeholder_text(slide, idx: int, text: str, bullets: Optional[list]):
    target = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            target = ph
            break
    if target is None:
        return
    tf = target.text_frame
    if bullets:
        tf.clear()
        for i, line in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = str(line)
            para.level = 0
    else:
        tf.text = str(text)


def build(plan: dict, out_path: Path, run_id: str = "") -> dict:
    master = plan.get("master")
    if master and Path(master).expanduser().exists():
        prs = Presentation(str(Path(master).expanduser().resolve()))
        # Remove any pre-existing slides from the template copy.
        xml_slides = prs.slides._sldIdLst
        for sld in list(xml_slides):
            xml_slides.remove(sld)
    else:
        prs = Presentation()  # 4:3 default; widescreen comes from master normally

    slides_spec = plan.get("slides", [])
    for s_i, spec in enumerate(slides_spec):
        layout = _resolve_layout(prs, spec.get("layout", 1))
        slide = prs.slides.add_slide(layout)
        placeholders = spec.get("placeholders", {})
        bullets = spec.get("bullets", {})
        for idx_str, text in placeholders.items():
            try:
                idx = int(idx_str)
            except (TypeError, ValueError):
                continue
            _set_placeholder_text(slide, idx, text, bullets.get(idx_str))
        notes = spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)
        if run_id:
            emit(run_id, "log", agent="builder", level="debug",
                 message=f"Built slide {s_i + 1}/{len(slides_spec)} "
                         f"using layout '{layout.name}'")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return {"slides": len(slides_spec), "out": str(out_path)}


def main() -> int:
    p = argparse.ArgumentParser(description="Build .pptx from a slide-plan JSON")
    p.add_argument("--plan", required=True)
    p.add_argument("--out", required=True)
    p.add_argument("--run", default="")
    args = p.parse_args()

    plan_path = Path(args.plan).expanduser().resolve()
    if not plan_path.exists():
        print(f"Plan not found: {plan_path}", file=sys.stderr)
        return 2
    plan = json.loads(plan_path.read_text(encoding="utf-8"))

    if args.run:
        emit(args.run, "tool_call", agent="builder", message="build_pptx",
             data={"plan": plan_path.name, "doc_type": plan.get("doc_type")})

    out_path = Path(args.out).expanduser().resolve()
    result = build(plan, out_path, run_id=args.run)

    if args.run:
        emit(args.run, "artifact", agent="builder",
             message=f"Generated {result['slides']} slide(s)",
             data={"out": result["out"], "kind": "pptx"})

    print(json.dumps({"ok": True, **result}))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
