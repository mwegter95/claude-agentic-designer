#!/usr/bin/env python3
"""evaluate.py — Deterministic rubric checks for design fidelity.

Tier-1 of the two-tier evaluation: fast, free, high-precision checks that the
Evaluator agent runs against the slide plan + design tokens + built .pptx. The
subjective Tier-2 criteria (visual hierarchy, narrative) are scored by Claude and
can be merged in via --llm-scores.

A hard violation of any deterministic criterion fails the gate regardless of the
LLM's aesthetic score.

Usage:
  python tools/evaluate.py --rubric skill/rubrics/slides.rubric.json \
      --plan run/slide_plan.json --tokens run/design_tokens.json \
      --pptx run/output.pptx --out run/evaluation.json \
      [--llm-scores run/llm_scores.json] [--run RUN123]
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Dict, List, Optional

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx import Presentation  # noqa: E402

from tools.common.events import emit  # noqa: E402


def _collect_pptx_facts(pptx_path: Path) -> Dict:
    prs = Presentation(str(pptx_path))
    fonts = set()
    colors = set()
    layouts_used = set()
    bullet_counts: List[int] = []
    for slide in prs.slides:
        layouts_used.add(slide.slide_layout.name)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            n_bullets = 0
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    n_bullets += 1
                for run in para.runs:
                    if run.font.name:
                        fonts.add(run.font.name)
                    if run.font.color and run.font.color.type is not None:
                        try:
                            colors.add("#" + str(run.font.color.rgb))
                        except Exception:
                            pass
            if n_bullets:
                bullet_counts.append(n_bullets)
    return {
        "fonts": sorted(fonts),
        "colors": sorted(colors),
        "layouts_used": sorted(layouts_used),
        "max_bullets": max(bullet_counts) if bullet_counts else 0,
        "slide_count": len(prs.slides._sldIdLst),
    }


def _check(cid: str, ok: bool, score: int, evidence: str, weight: float,
           fix: str = "") -> Dict:
    return {
        "criterion_id": cid,
        "tier": "deterministic",
        "pass": ok,
        "score": score,
        "weight": weight,
        "evidence": evidence,
        "fix_suggestion": fix if not ok else "",
    }


def run_deterministic(rubric: Dict, plan: Dict, tokens: Dict, facts: Dict) -> List[Dict]:
    findings: List[Dict] = []
    allowed_fonts = set(tokens.get("fonts", {}).get("allowed", []))
    allowed_hex = set(h.upper() for h in tokens.get("palette", {}).get("allowed_hex", []))
    allowed_layouts = {l["name"] for l in tokens.get("layouts", [])}
    rules = tokens.get("rules", {})

    by_id = {c["id"]: c for c in rubric.get("criteria", [])}

    if "fonts_on_brand" in by_id and allowed_fonts:
        off = [f for f in facts["fonts"] if f and f not in allowed_fonts]
        w = by_id["fonts_on_brand"]["weight"]
        findings.append(_check(
            "fonts_on_brand", not off, 5 if not off else 1,
            f"Off-brand fonts: {off}" if off else "All fonts within approved set",
            w, "Restrict runs to the master's major/minor typefaces."))

    if "colors_on_brand" in by_id and allowed_hex:
        off = [c for c in facts["colors"] if c.upper() not in allowed_hex]
        w = by_id["colors_on_brand"]["weight"]
        findings.append(_check(
            "colors_on_brand", not off, 5 if not off else 1,
            f"Off-palette colors: {off}" if off else "All colors within theme palette",
            w, "Use only theme color values from design_tokens.palette."))

    if "layouts_from_master" in by_id and allowed_layouts:
        off = [l for l in facts["layouts_used"] if l not in allowed_layouts]
        w = by_id["layouts_from_master"]["weight"]
        findings.append(_check(
            "layouts_from_master", not off, 5 if not off else 2,
            f"Non-master layouts: {off}" if off else "All layouts from master library",
            w, "Map every block to a named master layout."))

    if "bullet_density" in by_id:
        max_allowed = rules.get("max_bullets_per_slide", 6)
        ok = facts["max_bullets"] <= max_allowed
        w = by_id["bullet_density"]["weight"]
        findings.append(_check(
            "bullet_density", ok, 5 if ok else 2,
            f"Max bullets/slide = {facts['max_bullets']} (limit {max_allowed})",
            w, f"Split slides exceeding {max_allowed} bullets."))

    return findings


def merge_llm(findings: List[Dict], llm_scores: Optional[Dict], rubric: Dict) -> List[Dict]:
    if not llm_scores:
        return findings
    by_id = {c["id"]: c for c in rubric.get("criteria", [])}
    for item in llm_scores.get("scores", []):
        cid = item.get("criterion_id")
        weight = by_id.get(cid, {}).get("weight", 1.0)
        score = int(item.get("score", 3))
        threshold = by_id.get(cid, {}).get("pass_threshold", 3)
        findings.append({
            "criterion_id": cid,
            "tier": "llm",
            "pass": score >= threshold,
            "score": score,
            "weight": weight,
            "evidence": item.get("evidence", ""),
            "fix_suggestion": item.get("fix_suggestion", ""),
        })
    return findings


def aggregate(findings: List[Dict], rubric: Dict) -> Dict:
    total_w = sum(f["weight"] for f in findings) or 1.0
    weighted = sum((f["score"] / 5.0) * f["weight"] for f in findings) / total_w
    hard_fail = any(
        (not f["pass"]) and f["tier"] == "deterministic" for f in findings
    )
    threshold = rubric.get("pass_threshold", 0.8)
    passed = (weighted >= threshold) and not hard_fail
    return {
        "weighted_score": round(weighted, 4),
        "threshold": threshold,
        "hard_fail": hard_fail,
        "passed": passed,
    }


def main() -> int:
    p = argparse.ArgumentParser(description="Run rubric evaluation")
    p.add_argument("--rubric", required=True)
    p.add_argument("--plan", required=True)
    p.add_argument("--tokens", required=True)
    p.add_argument("--pptx", required=True)
    p.add_argument("--out", required=True)
    p.add_argument("--llm-scores", default="")
    p.add_argument("--run", default="")
    args = p.parse_args()

    rubric = json.loads(Path(args.rubric).read_text(encoding="utf-8"))
    plan = json.loads(Path(args.plan).read_text(encoding="utf-8"))
    tokens = json.loads(Path(args.tokens).read_text(encoding="utf-8"))
    facts = _collect_pptx_facts(Path(args.pptx).expanduser().resolve())

    if args.run:
        emit(args.run, "tool_call", agent="evaluator", message="evaluate",
             data={"rubric": rubric.get("name")})

    findings = run_deterministic(rubric, plan, tokens, facts)
    llm_scores = None
    if args.llm_scores and Path(args.llm_scores).exists():
        llm_scores = json.loads(Path(args.llm_scores).read_text(encoding="utf-8"))
    findings = merge_llm(findings, llm_scores, rubric)
    summary = aggregate(findings, rubric)

    report = {
        "rubric": rubric.get("name"),
        "facts": facts,
        "findings": findings,
        "summary": summary,
    }
    Path(args.out).write_text(json.dumps(report, indent=2), encoding="utf-8")

    if args.run:
        emit(args.run, "evaluation", agent="evaluator",
             level="info" if summary["passed"] else "warn",
             message=f"Score {summary['weighted_score']:.2f} / "
                     f"threshold {summary['threshold']} — "
                     f"{'PASS' if summary['passed'] else 'REVISE'}",
             data=summary)

    print(json.dumps({"ok": True, **summary}))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
